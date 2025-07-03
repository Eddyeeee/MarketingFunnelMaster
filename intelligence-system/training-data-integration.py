"""
Training Data Integration System
Processes sales materials from das-magische-tool and other sources
Integrates with AI Research Engine for pattern recognition
"""

import os
import json
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional
import pandas as pd
from datetime import datetime

# Document processing libraries
import PyPDF2
from pptx import Presentation
from bs4 import BeautifulSoup
import re

class TrainingDataProcessor:
    """
    Processes training materials and extracts patterns for AI learning
    """
    
    def __init__(self, base_path: str = "../training-data"):
        self.base_path = Path(base_path)
        self.processed_data = {
            "sales_patterns": [],
            "funnel_structures": [],
            "design_elements": [],
            "conversion_triggers": [],
            "content_templates": []
        }
        
        # Setup logging
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)
        
    def process_das_magische_tool_materials(self) -> Dict[str, Any]:
        """
        Process all materials from das-magische-tool directory
        """
        tool_path = self.base_path / "das-magische-tool"
        
        if not tool_path.exists():
            self.logger.warning(f"Das-magische-tool directory not found: {tool_path}")
            return {}
            
        results = {
            "folien_patterns": self._process_folien(tool_path / "folien"),
            "funnel_patterns": self._process_funnelseiten(tool_path / "funnelseite"),  
            "additional_materials": self._process_sonstiges(tool_path / "sonstiges")
        }
        
        return results
    
    def _process_folien(self, folien_path: Path) -> List[Dict]:
        """
        Extract patterns from presentation slides
        """
        patterns = []
        
        if not folien_path.exists():
            return patterns
            
        for file in folien_path.glob("*"):
            if file.suffix.lower() == '.pdf':
                patterns.extend(self._extract_pdf_patterns(file))
            elif file.suffix.lower() in ['.ppt', '.pptx']:
                patterns.extend(self._extract_pptx_patterns(file))
                
        return patterns
    
    def _process_funnelseiten(self, funnel_path: Path) -> List[Dict]:
        """
        Extract funnel structures and conversion elements
        """
        funnel_patterns = []
        
        if not funnel_path.exists():
            return funnel_patterns
            
        for file in funnel_path.glob("*"):
            if file.suffix.lower() in ['.html', '.htm']:
                funnel_patterns.extend(self._extract_html_patterns(file))
            elif file.suffix.lower() == '.pdf':
                funnel_patterns.extend(self._extract_pdf_patterns(file))
                
        return funnel_patterns
    
    def _process_sonstiges(self, sonstiges_path: Path) -> List[Dict]:
        """
        Process additional materials
        """
        additional = []
        
        if not sonstiges_path.exists():
            return additional
            
        for file in sonstiges_path.rglob("*"):
            if file.is_file():
                additional.append({
                    "file_type": file.suffix,
                    "file_name": file.name,
                    "processed_date": datetime.now().isoformat(),
                    "content_type": self._identify_content_type(file)
                })
                
        return additional
    
    def _extract_pdf_patterns(self, pdf_path: Path) -> List[Dict]:
        """
        Extract text patterns from PDF files
        """
        patterns = []
        
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = ""
                
                for page in pdf_reader.pages:
                    text_content += page.extract_text()
                
                # Extract key patterns
                patterns.append({
                    "source_file": pdf_path.name,
                    "content_type": "pdf_text",
                    "extracted_text": text_content,
                    "headlines": self._extract_headlines(text_content),
                    "calls_to_action": self._extract_ctas(text_content),
                    "benefits": self._extract_benefits(text_content),
                    "processed_date": datetime.now().isoformat()
                })
                
        except Exception as e:
            self.logger.error(f"Error processing PDF {pdf_path}: {e}")
            
        return patterns
    
    def _extract_pptx_patterns(self, pptx_path: Path) -> List[Dict]:
        """
        Extract patterns from PowerPoint presentations
        """
        patterns = []
        
        try:
            prs = Presentation(pptx_path)
            slide_contents = []
            
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                        
                slide_contents.append(slide_text.strip())
            
            patterns.append({
                "source_file": pptx_path.name,
                "content_type": "presentation",
                "slide_count": len(slide_contents),
                "slide_contents": slide_contents,
                "key_messages": self._extract_key_messages(slide_contents),
                "processed_date": datetime.now().isoformat()
            })
            
        except Exception as e:
            self.logger.error(f"Error processing PPTX {pptx_path}: {e}")
            
        return patterns
    
    def _extract_html_patterns(self, html_path: Path) -> List[Dict]:
        """
        Extract funnel patterns from HTML files
        """
        patterns = []
        
        try:
            with open(html_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
                soup = BeautifulSoup(html_content, 'html.parser')
                
                patterns.append({
                    "source_file": html_path.name,
                    "content_type": "html_funnel",
                    "page_title": soup.find('title').text if soup.find('title') else "",
                    "headlines": [h.text for h in soup.find_all(['h1', 'h2', 'h3'])],
                    "buttons": [btn.text for btn in soup.find_all(['button', 'a']) if btn.text.strip()],
                    "forms": len(soup.find_all('form')),
                    "images": len(soup.find_all('img')),
                    "videos": len(soup.find_all(['video', 'iframe'])),
                    "processed_date": datetime.now().isoformat()
                })
                
        except Exception as e:
            self.logger.error(f"Error processing HTML {html_path}: {e}")
            
        return patterns
    
    def _extract_headlines(self, text: str) -> List[str]:
        """Extract potential headlines from text"""
        lines = text.split('\n')
        headlines = []
        
        for line in lines:
            line = line.strip()
            # Look for short, impactful lines that could be headlines
            if 5 <= len(line) <= 100 and line.isupper() or '!' in line:
                headlines.append(line)
                
        return headlines[:10]  # Return top 10 potential headlines
    
    def _extract_ctas(self, text: str) -> List[str]:
        """Extract Call-to-Action phrases"""
        cta_patterns = [
            r'Jetzt\s+\w+',
            r'Hier\s+klicken',
            r'Kostenlos\s+\w+',
            r'Sofort\s+\w+',
            r'Jetzt\s+starten',
            r'Mehr\s+erfahren'
        ]
        
        ctas = []
        for pattern in cta_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            ctas.extend(matches)
            
        return list(set(ctas))
    
    def _extract_benefits(self, text: str) -> List[str]:
        """Extract benefit statements"""
        benefit_indicators = ['vorteile', 'nutzen', 'profitieren', 'sparen', 'verdienen']
        benefits = []
        
        sentences = text.split('.')
        for sentence in sentences:
            for indicator in benefit_indicators:
                if indicator in sentence.lower():
                    benefits.append(sentence.strip())
                    break
                    
        return benefits[:5]  # Return top 5 benefits
    
    def _extract_key_messages(self, slide_contents: List[str]) -> List[str]:
        """Extract key messages from presentation slides"""
        all_text = ' '.join(slide_contents)
        
        # Simple extraction of key messages (can be enhanced with NLP)
        key_messages = []
        for content in slide_contents:
            if len(content.strip()) > 10 and len(content.strip()) < 200:
                key_messages.append(content.strip())
                
        return key_messages
    
    def _identify_content_type(self, file_path: Path) -> str:
        """Identify the type of content based on file extension"""
        extension_map = {
            '.pdf': 'document',
            '.pptx': 'presentation',
            '.ppt': 'presentation',
            '.html': 'webpage',
            '.htm': 'webpage',
            '.docx': 'document',
            '.doc': 'document',
            '.txt': 'text',
            '.jpg': 'image',
            '.png': 'image',
            '.gif': 'image',
            '.mp4': 'video',
            '.avi': 'video'
        }
        
        return extension_map.get(file_path.suffix.lower(), 'unknown')
    
    def export_training_data(self, output_path: str = "../intelligence-system/research-data/training-patterns.json"):
        """
        Export processed training data for AI system integration
        """
        processed_materials = self.process_das_magische_tool_materials()
        
        training_export = {
            "export_date": datetime.now().isoformat(),
            "source_materials": "das-magische-tool",
            "processed_patterns": processed_materials,
            "pattern_count": len(str(processed_materials)),
            "integration_ready": True
        }
        
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(training_export, f, indent=2, ensure_ascii=False)
            
        self.logger.info(f"Training data exported to: {output_file}")
        return output_file

if __name__ == "__main__":
    # Process training materials
    processor = TrainingDataProcessor()
    
    # Export for AI integration
    export_path = processor.export_training_data()
    print(f"Training data ready for AI integration: {export_path}")